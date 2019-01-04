from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os
import sys
sys.path.append(os.getcwd())

#PATH_TO_PRESENTATION = 'test.pptx'


def change_color(shape):
	# Color
	fill = shape.fill
	fill.solid()
	# RGB (255,255,218) --> Coverage Area
	fill.fore_color.rgb = RGBColor(255, 255, 218)

	# Line
	line = shape.line
	line.color.rgb = RGBColor(255, 0, 0)
	line.dash_style = MSO_LINE_DASH_STYLE.DASH

	# shadow
	shape.shadow.inherit = False

	return shape


prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes

# shapes.title.text = 'Adding an AutoShape'
left = Inches(4.0) # 0.93" centers this overall set of shapes
top = Inches(4.0)
width = Inches(0.75)
height = Inches(0.5)

#shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)

#left = left + width - Inches(0.2)
#width = Inches(0.5) # chevrons need more width for visual balance

for n in range(1, 10):
	shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
	shape = change_color(shape)
	left = left + width - Inches(0.2)

prs.save('test.pptx')
