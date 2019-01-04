from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import random
import os
import sys
sys.path.append(os.getcwd())

#PATH_TO_PRESENTATION = 'test.pptx'

# Asking
# What is the magnification of the map?
# Ceiling Height?
# Which Sensor? (Aurora/BS/Axis/V2)
# Number?


def change_color(shape):
	# Color
	fill = shape.fill
	fill.solid()

	# RGB (255,255,218) --> Coverage Area
	r, g, b = random.randint(0, 255), random.randint(0, 255), random.randint(0, 255)
	fill.fore_color.rgb = RGBColor(r, g, b)
	# sets opacity to 80%
	# fill.transparency = 0.25
	# spPr = fill._xPr
	# print(spPr.xml)

	# Line
	line = shape.line
	line.color.rgb = RGBColor(255, 0, 0)
	line.dash_style = MSO_LINE_DASH_STYLE.DASH

	return shape

# answer = input("Enter ")
# print(answer)

prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes

# shapes.title.text = 'Adding an AutoShape'
left = Inches(4.0) # 0.93" centers this overall set of shapes
top = Inches(4.0)
width = Inches(0.2)
height = Inches(0.2)

shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)

left = left + width - Inches(0.05)
width = Inches(0.2) # chevrons need more width for visual balance

for n in range(2, 4):
	shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
	shape = change_color(shape)
	left = left + width - Inches(0.01)

prs.save('test2.pptx')
